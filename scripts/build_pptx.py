#!/usr/bin/env python3
"""Generate docs/Agora-pitch.pptx — a dark, on-brand pitch deck.

Mirrors docs/SLIDES.md / docs/pitch.html. Run: python3 scripts/build_pptx.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

INK = RGBColor(0x0A, 0x0E, 0x1A)
INK2 = RGBColor(0x12, 0x18, 0x2B)
WHITE = RGBColor(0xE6, 0xE9, 0xF2)
GRAY = RGBColor(0x8A, 0x90, 0xA6)
PLANNER = RGBColor(0x7C, 0x83, 0xFF)
TUTOR = RGBColor(0x3E, 0xCF, 0x8E)
ASSESSOR = RGBColor(0xFF, 0xB0, 0x20)
DIAG = RGBColor(0xFF, 0x5D, 0x73)
LEARNER = RGBColor(0x22, 0xD3, 0xEE)
ORCH = RGBColor(0xC0, 0x84, 0xFC)

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(1, 0, 0, EMU_W, EMU_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = INK
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def para(tf, runs, size, bold=False, space_after=6, align=PP_ALIGN.LEFT, line=1.1, first=False):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_after = Pt(space_after)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if isinstance(runs, str):
        runs = [(runs, WHITE)]
    for text, color in runs:
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        r.font.name = "Arial"
    return p


def kicker(s, text):
    tf = textbox(s, 0.9, 0.7, 11, 0.6)
    para(tf, [(text.upper(), GRAY)], 13, bold=True, first=True)


def title(s, runs, y=1.4, size=40):
    tf = textbox(s, 0.9, y, 11.5, 2.2)
    para(tf, runs, size, bold=True, line=1.02, first=True)


def lead(s, text, y=3.3):
    tf = textbox(s, 0.9, y, 10.8, 2.5)
    para(tf, [(text, RGBColor(0xB8, 0xBE, 0xCE))], 20, line=1.35, first=True)


# 1 — Title
s = slide()
tf = textbox(s, 0.9, 0.9, 11, 0.8)
para(tf, [("\U0001F3DB  Agora", WHITE)], 26, bold=True, first=True)
para(tf, [("The self-driving classroom", GRAY)], 14)
title(s, [("Autonomous AI agents that ", WHITE), ("teach a goal end-to-end.", TUTOR)], y=2.6, size=44)
lead(s, "A multi-agent tutoring system for the Agentic AI in Education hackathon. "
        "Give it a learning goal; six agents plan, teach, diagnose, and adapt — "
        "with little to no human intervention.", y=4.6)

# 2 — Problem
s = slide()
kicker(s, "The problem")
title(s, [("The best teaching method doesn't scale.", WHITE)], size=38)
lead(s, "One-to-one tutoring produces a two-standard-deviation improvement (Bloom, 1984): "
        "the most effective intervention we have in education — and the least scalable. "
        "There will never be enough expert tutors.")
tf = textbox(s, 0.9, 5.4, 11, 1)
para(tf, [("+2σ  ", PLANNER), ("one-to-one tutoring    vs    one teacher, thirty students", GRAY)], 22, bold=True, first=True)

# 3 — Insight
s = slide()
kicker(s, "The insight")
title(s, [("Tutoring is a ", WHITE), ("loop", LEARNER), (", not a Q&A.", WHITE)], size=38)
lead(s, "Most “AI tutors” just answer questions. A real tutor plans what to teach, "
        "notices confusion, diagnoses the underlying gap, and changes the plan to fix it.")
tf = textbox(s, 0.9, 5.3, 11.5, 1)
para(tf, [("plan  →  teach  →  assess  →  ", WHITE), ("diagnose", DIAG),
          ("  →  adapt  ↺", WHITE)], 22, bold=True, first=True)

# 4 — Solution
s = slide()
kicker(s, "The solution")
title(s, [("Six autonomous agents run the whole loop.", WHITE)], size=36)
lead(s, "One input — a learning goal. Everything else is decided by the system. "
        "In autopilot, a simulated learner even answers the quizzes, so a full session "
        "runs with zero human input.")

# 5 — Crew
s = slide()
kicker(s, "The crew")
title(s, [("Each agent, one job.", WHITE)], size=36)
agents = [
    ("◎ Orchestrator", "Decides the next move", ORCH),
    ("✦ Planner", "Designs the learning path", PLANNER),
    ("❂ Tutor", "Teaches each concept", TUTOR),
    ("◈ Assessor", "Checks understanding", ASSESSOR),
    ("⚕ Diagnostician", "Finds gaps, triggers remediation", DIAG),
    ("◍ Learner", "Answers and progresses", LEARNER),
]
for i, (name, role, color) in enumerate(agents):
    col = i % 3
    row = i // 3
    x = 0.9 + col * 4.0
    y = 2.5 + row * 1.7
    tf = textbox(s, x, y, 3.8, 1.5)
    para(tf, [(name, color)], 20, bold=True, first=True)
    para(tf, [(role, GRAY)], 14)

# 6 — Emergent remediation
s = slide()
kicker(s, "The headline behavior")
title(s, [("It ", WHITE), ("rewrites its own plan.", PLANNER)], size=38)
lead(s, "When a wrong answer traces to a missing prerequisite, the Orchestrator injects a "
        "remedial micro-lesson it never originally planned, teaches it, then returns to the "
        "blocked concept. Nobody scripted it — it's a decision from the learner's state.")
tf = textbox(s, 0.9, 5.4, 11.6, 1)
para(tf, [("missed check  →  ", WHITE), ("misconception found", DIAG),
          ("  →  ", WHITE), ("plan rewritten", PLANNER), ("  →  gap closed", WHITE)], 18, bold=True, first=True)

# 7 — Demo
s = slide()
kicker(s, "Live demo")
title(s, [("Watch it run.", WHITE)], size=40)
lead(s, "Autopilot run of “Understand recursion”: the plan builds, agents light up as they "
        "act, the activity stream shows every autonomous decision, and a remediation lesson is "
        "auto-added live. Then an interactive run with a human answering.")
tf = textbox(s, 0.9, 5.6, 11, 0.6)
para(tf, [("→ switch to the app", GRAY)], 16, first=True)

# 8 — Engineering
s = slide()
kicker(s, "Engineering")
title(s, [("Production-shaped, demo-safe.", WHITE)], size=36)
left = textbox(s, 0.9, 2.6, 5.6, 4)
para(left, [("Stack", TUTOR)], 18, bold=True, first=True)
for t in ["Node + Express (tsx) · React control room",
          "Redis: state, replayable events, pub/sub, analytics",
          "Real-time via SSE + Redis pub/sub — scales out",
          "One-command deploy to Scalingo (+ Redis addon)"]:
    para(left, [("• " + t, RGBColor(0xB8, 0xBE, 0xCE))], 15, line=1.2)
right = textbox(s, 6.9, 2.6, 5.6, 4)
para(right, [("Why it never breaks on stage", ASSESSOR)], 18, bold=True, first=True)
for t in ["Provider-agnostic LLM (Anthropic / OpenAI)…",
          "…that degrades to a deterministic engine",
          "No API key, no network needed for the demo",
          "Typed · unit + integration + smoke + browser tested"]:
    para(right, [("• " + t, RGBColor(0xB8, 0xBE, 0xCE))], 15, line=1.2)

# 9 — Why it scores
s = slide()
kicker(s, "Why it scores")
title(s, [("Built against the rubric.", WHITE)], size=36)
crit = [
    ("Autonomy", "Closed multi-step loop to a goal, no human needed", LEARNER),
    ("Innovation", "Multi-agent orchestration + self-rewriting plans", PLANNER),
    ("Impact", "Scalable 1:1 tutoring for every learner", TUTOR),
    ("Theme", "Agentic by design, with emergent behavior", ORCH),
    ("Quality", "Typed, tested, verified, polished demo", ASSESSOR),
]
for i, (t, d, c) in enumerate(crit):
    x = 0.9 + (i % 3) * 4.0
    y = 2.7 + (i // 3) * 1.9
    tf = textbox(s, x, y, 3.8, 1.7)
    para(tf, [(t, c)], 18, bold=True, first=True)
    para(tf, [(d, GRAY)], 13, line=1.2)

# 10 — Close
s = slide()
tf = textbox(s, 0.9, 0.9, 11, 0.8)
para(tf, [("\U0001F3DB  Agora", WHITE)], 24, bold=True, first=True)
title(s, [("Autonomous, adaptive tutoring ", WHITE), ("for every learner.", TUTOR)], y=2.6, size=38)
lead(s, "And fittingly — this project about agentic AI was built by an agentic workflow: "
        "agents writing code, reviewing each other, and verifying the result before shipping.",
     y=4.6)
tf = textbox(s, 0.9, 6.2, 11, 0.6)
para(tf, [("github.com/iliaskalalou/agora-agentic-tutor", GRAY)], 15, first=True)

out = os.path.join(os.path.dirname(__file__), "..", "docs", "Agora-pitch.pptx")
prs.save(out)
print("wrote", os.path.normpath(out), "with", len(prs.slides._sldIdLst), "slides")
